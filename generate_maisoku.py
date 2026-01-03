from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image
import os

# Configuration
WIDTH = Cm(29.7)
HEIGHT = Cm(21.0)
MARGIN = Cm(1.0)
COLOR_TEXT = RGBColor(74, 64, 54) # #4A4036
COLOR_BG = RGBColor(253, 251, 247) # #FDFBF7
COLOR_ACCENT = RGBColor(157, 175, 115) # #9DAF73
COLOR_HEADING = RGBColor(62, 58, 57) # #3E3A39

def ensure_image(path):
    """Converts image to compatible format if necessary."""
    try:
        img = Image.open(path)
        if img.format not in ['JPEG', 'PNG', 'BMP', 'GIF', 'TIFF']:
            # Convert to RGB (remove alpha if present for JPEG) and save as JPEG
            if img.mode in ("RGBA", "P"):
                img = img.convert("RGB")
            new_path = path + ".converted.jpg"
            img.save(new_path, "JPEG")
            return new_path
        return path
    except Exception as e:
        print(f"Error processing image {path}: {e}")
        return path

def create_presentation():
    prs = Presentation()
    prs.slide_width = WIDTH
    prs.slide_height = HEIGHT
    
    # Create blank slide layout
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Background - Tiled Pattern
    try:
        bg_pattern_path = "style/bg_pattern.png"
        bg_pattern = Image.open(bg_pattern_path)
        
        # Create a new image of slide size
        # Convert Cm to pixels (approx 96 DPI for calculation, but PIL uses pixels)
        # 1 Cm = ~37.8 pixels
        dpi = 96
        width_px = int(WIDTH.inches * dpi)
        height_px = int(HEIGHT.inches * dpi)
        
        tiled_bg = Image.new('RGB', (width_px, height_px))
        
        # Tile the pattern
        for x in range(0, width_px, bg_pattern.width):
            for y in range(0, height_px, bg_pattern.height):
                tiled_bg.paste(bg_pattern, (x, y))
        
        # Save temp background
        bg_temp_path = "temp_bg_tiled.jpg"
        tiled_bg.save(bg_temp_path, "JPEG", quality=95)
        
        # Add as picture stretched to fill slide (it's already correct aspect ratio)
        slide.shapes.add_picture(bg_temp_path, 0, 0, width=WIDTH, height=HEIGHT)
        
        # Send to back? PPTX adds in order, so adding first effectively puts it at back.
        
    except Exception as e:
        print(f"Failed to create tiled background: {e}")
        # Fallback to solid color
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = COLOR_BG

    # --- Header Row (Y=1.0 - 4.5) ---
    y_header = MARGIN
    h_header = Cm(3.5)
    
    # Catchphrase Box
    left_catch = MARGIN
    width_catch = Cm(11.0)
    txBox = slide.shapes.add_textbox(left_catch, y_header, width_catch, h_header)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "都心近接の時短邸宅"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COLOR_HEADING
    p.font.name = "Zen Old Mincho"
    
    p = tf.add_paragraph()
    p.text = "快適な子育てを応援するZEH住宅！\n新宿20分！保育施設/小学校/球技OK公園5分圏内！家事代行級の設備！"
    p.font.size = Pt(11)
    p.font.color.rgb = COLOR_TEXT
    p.font.name = "Noto Sans JP"

    # QR Box
    left_qr = left_catch + width_catch + Cm(0.5)
    width_qr = Cm(5.0)
    # Image
    img_path = ensure_image("QRcode/QR_Only.png")
    slide.shapes.add_picture(img_path, left_qr + Cm(1.0), y_header, height=Cm(2.5))
    # Text
    txBox = slide.shapes.add_textbox(left_qr, y_header + Cm(2.6), width_qr, Cm(1.0))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "詳細情報はこちら！"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(10)
    p.font.bold = True

    # Title Box
    left_title = left_qr + width_qr + Cm(0.5)
    width_title = Cm(10.7)
    # Background rect
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_title, y_header, width_title, h_header)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_HEADING
    shape.line.fill.background() # No border
    
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = "Maison Terrasse Asagaya"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    p = tf.add_paragraph()
    p.text = "3LDK 81.18㎡"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # --- Main Grid (Y=5.0 - 16.0) ---
    y_main = Cm(5.0)
    h_main = Cm(11.5)
    
    # Left Col
    x_col1 = MARGIN
    w_col1 = Cm(9.0)
    
    # Points
    points = [
        ("1. 駅徒歩8分", "都心アクセスと閑静さを両立！"),
        ("2. 徒歩5分圏内に子育て施設充実", "保育園/幼稚園/小学校/球技OK公園"),
    ]

    current_y = y_main
    
    # Draw Points
    for title, desc in points:
        height_approx = Cm(1.8)
        txBox = slide.shapes.add_textbox(x_col1, current_y, w_col1, height_approx)
        tf = txBox.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.bold = True
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(184, 134, 11) # Gold
        
        p = tf.add_paragraph()
        p.text = desc
        p.font.size = Pt(10)
        
        current_y += height_approx
    
    # --- Premium Equipment Box ---
    # Space below points
    current_y += Cm(0.5)
    
    # Calculate remaining height or set fixed
    box_height = Cm(8.0)
    
    # Blue Box Background
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_col1, current_y, w_col1, box_height)
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(20, 40, 80) # Navy Blue
    box.line.color.rgb = RGBColor(13, 27, 54)
    
    # Inside Box - Y Coordinates relative to slide, but logically inside
    inner_y = current_y + Cm(0.3)
    inner_x = x_col1 + Cm(0.3)
    inner_w = w_col1 - Cm(0.6)
    
    # Title
    txBox = slide.shapes.add_textbox(inner_x, inner_y, inner_w, Cm(1.0))
    p = txBox.text_frame.paragraphs[0]
    p.text = "家事代行級の設備"
    p.font.name = "Zen Old Mincho"
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(255, 224, 130) # Gold
    p.alignment = PP_ALIGN.CENTER
    
    inner_y += Cm(1.2)
    
    # Row 1: Bath
    # Text
    txBox = slide.shapes.add_textbox(inner_x, inner_y, inner_w - Cm(2.8), Cm(2.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "自動おそうじ浴槽"
    p.font.bold = True
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(255, 224, 130)
    p = tf.add_paragraph()
    p.text = "洗剤を入れてスイッチを押すだけで浴槽をキレイに！"
    p.font.size = Pt(8)
    p.font.color.rgb = RGBColor(255, 255, 255)
    # Image
    slide.shapes.add_picture(ensure_image("images/101/bath.JPG"), inner_x + inner_w - Cm(2.5), inner_y, width=Cm(2.5), height=Cm(1.8))
    
    inner_y += Cm(2.0)
    
    # Row 2: Kitchen
    # Text
    txBox = slide.shapes.add_textbox(inner_x, inner_y, inner_w - Cm(2.8), Cm(2.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "食洗機付きキッチン"
    p.font.bold = True
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(255, 224, 130)
    p = tf.add_paragraph()
    p.text = "リビングを見渡しやすいペニンシュラキッチン！"
    p.font.size = Pt(8)
    p.font.color.rgb = RGBColor(255, 255, 255)
    # Image
    slide.shapes.add_picture(ensure_image("images/101/kitchen.JPG"), inner_x + inner_w - Cm(2.5), inner_y, width=Cm(2.5), height=Cm(1.8))
    
    inner_y += Cm(2.2)
    
    # Summary
    txBox = slide.shapes.add_textbox(inner_x, inner_y, inner_w, Cm(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "1ヶ月で約26時間分の家事を削減！\n約32,000円分の時間価値！"
    p.font.bold = True
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Center Col
    x_col2 = x_col1 + w_col1 + Cm(0.5)
    w_col2 = Cm(10.0)
    
    # Info List
    txBox = slide.shapes.add_textbox(x_col2, y_main, w_col2, Cm(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "✓ 各部屋の写真    ✓ 周辺地図    ✓ 内見ガイド"
    p.font.size = Pt(11)
    p.alignment = PP_ALIGN.CENTER
    
    # Plan
    slide.shapes.add_picture(ensure_image("images/101/plan_101.png"), x_col2, y_main + Cm(2.0), width=w_col2) # Approximate fit

    # Right Col
    x_col3 = x_col2 + w_col2 + Cm(0.5)
    w_col3 = Cm(7.7)
    
    # Conditions Box
    h_cond = Cm(6.0)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x_col3, y_main, w_col3, h_cond)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_ACCENT
    # shape.fill.fore_color.brightness = 0.8 # Lighten not easily supported in simple API, use color matching
    shape.fill.fore_color.rgb = RGBColor(244, 247, 239) # Light accent
    shape.line.color.rgb = COLOR_TEXT
    
    tf = shape.text_frame
    tf.margin_top = Cm(0.5)
    tf.margin_left = Cm(0.5)
    p = tf.paragraphs[0]
    p.text = "賃料など条件"
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    conditions = ["賃料: ", "管理費: ", "敷金/礼金: ", "入居時期: "]
    for c in conditions:
        p = tf.add_paragraph()
        p.text = c
        p.space_before = Pt(14)
        # Add underline effect manually if needed, or leave blank for filling

    # Deck Image
    slide.shapes.add_picture(ensure_image("images/101/deck1.JPG"), x_col3, y_main + h_cond + Cm(0.5), width=w_col3, height=Cm(5.0))
    # Label
    txBox = slide.shapes.add_textbox(x_col3, y_main + h_cond + Cm(0.5) + Cm(5.0) - Cm(0.8), w_col3, Cm(0.8))
    p = txBox.text_frame.paragraphs[0]
    p.text = "ウッドデッキ"
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER


    # --- Footer (Y=17.0 - 20.0) ---
    y_footer = Cm(17.0)
    
    # Annotation
    txBox = slide.shapes.add_textbox(MARGIN, y_footer, Cm(20.0), Cm(1.0))
    p = txBox.text_frame.paragraphs[0]
    p.text = "※1 31日間毎日 浴槽掃除5分/食器洗い15分3回行った場合   ※2 東京都の最低賃金1,226円/時での概算"
    p.font.size = Pt(8)
    
    # Agent Box
    y_agent = y_footer + Cm(1.0)
    h_agent = Cm(2.0)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN, y_agent, WIDTH - MARGIN*2, h_agent)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
    shape.line.color.rgb = COLOR_HEADING
    shape.line.width = Pt(2)
    
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = "業者記入欄"
    p.font.color.rgb = RGBColor(200, 200, 200)
    p.alignment = PP_ALIGN.CENTER

    prs.save('MaisonTerrasseAsagaya_Maisoku_v4.pptx')
    print("PPTX generated successfully.")

if __name__ == "__main__":
    create_presentation()
